
import collections 
import collections.abc
from pptx import Presentation
import sys
import io

# 设置标准输出为 UTF-8 编码，防止中文乱码
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

def extract_content(file_path):
    prs = Presentation(file_path)
    content = []
    
    for i, slide in enumerate(prs.slides):
        content.append(f"## Slide {i+1}")
        
        # 提取标题
        if slide.shapes.title:
            content.append(f"### Title: {slide.shapes.title.text}")
        
        # 提取所有文本框中的内容
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # 排除已经是标题的内容
                if slide.shapes.title and shape == slide.shapes.title:
                    continue
                content.append(shape.text.strip())
        
        # 提取幻灯片备注
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                content.append(f"**Notes:** {notes}")
        
        content.append("-" * 20)
    
    return "\n\n".join(content)

if __name__ == "__main__":
    file_path = r"AI天气预测.pptx"
    extracted_text = extract_content(file_path)
    with open("extracted_content.txt", "w", encoding="utf-8") as f:
        f.write(extracted_text)
    print("Content extracted successfully to extracted_content.txt")
